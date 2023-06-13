import os
from PIL import Image
import json
import requests
import pptx
from pptx.util import Inches
import numpy as np
import csv
import sys
import cv2

from json import JSONEncoder
from selenium import webdriver
from webdriver_manager.chrome import ChromeDriverManager
from selenium.webdriver.common.by import By
from selenium.common.exceptions import NoSuchElementException
from flask import Flask
from tkinter import Tk

import tree_maker as tm
from bs4 import BeautifulSoup as BS
import re

app = Flask(__name__)
@app.route("/")
def index():
    return "Hello World!"

valDir = "values.json"
plantsList = "plants.csv"
id_dir = "id_dir"
searchUrl = "https://www.biolib.cz/cz/main/"

resultDict = {}
resultTree = []

class plantEncoder(JSONEncoder):
        def default(self, o):
            return {o.__class__.__name__: o.__dict__}

class plant:
    def __init__(self, kingdom, phylum, plantClass, order, family, genus, czechName, latinName):
        self.kingdom = kingdom
        self.phylum = phylum
        self.plantClass = plantClass
        self.order = order
        self.family = family
        self.genus = genus
        self.czechName = czechName
        self.latinName = latinName

def produce_output():
    imgs = os.listdir(id_dir)
    prs = pptx.Presentation()
    lyt=prs.slide_layouts[0] # choosing a slide layout

    for image in imgs:
        if image in [".DS_Store", ".gitkeep"]:
            continue

        imgPath = os.path.join(id_dir, image)
        # im = Image.open(imgPath)
        # rgb_im = im.convert("RGB")
        # rgb_im.save(f"{image[:image.find('.')]}.jpg")

        hierarchy = plants[image[:image.find('.')]]
        if hierarchy == None:
            continue

        slide=prs.slides.add_slide(lyt) #New slide
        title=slide.shapes.title
        title.top = Inches(0)
        title.left = Inches(4.5)
        title.width = Inches(4)
        title.height = Inches(4)
        print(imgPath)
        imag=slide.shapes.add_picture(imgPath, Inches(0.2), Inches(0.2), height=Inches(3.333)) #Image
        subtitle=slide.placeholders[1]

        plantName = imgPath[:imgPath.find('.')]
        title.text=f"{str(hierarchy.czechName).capitalize()} ({str(hierarchy.latinName).capitalize()})"
        subtitle.text = f"Rod: {str(hierarchy.genus).capitalize()}, \nDruh: {str(hierarchy.family).capitalize()}, \nČeleď: {str(hierarchy.order).capitalize()}"

        prs.save("plants.pptx")

plants = {}
#python main.py [bool train data + make powerpoint] [bool use list]
def main():
    if (len(sys.argv) in [2, 3] and strToBool(sys.argv[1])):
        eraseFile()
        driver = webdriver.Chrome(ChromeDriverManager().install())
        driver.get(searchUrl)

        consent = driver.find_element(By.ID, 'consentAllButton')
        consent.click()

        if strToBool(sys.argv[2]):
            with open(plantsList, "r") as f:
                csvreader = csv.reader(f)
                for i in csvreader:
                    for j in i:
                        plants[j] = getHierarchy(j, driver)
                        print(j)
        else:
            for j in os.listdir(id_dir):
                if j in [".DS_Store", ".gitkeep"]:
                    continue
                plants[j[:j.find(".")]] = getHierarchy(j[:j.find(".")], driver)
                print(j[:j.find(".")])

        driver.close()
        json_object = json.dumps(resultDict, indent=4)
        with open(valDir, "a") as f:
            f.write(json_object)

    #Whether or not to use the plants.csv file as the source (True: yes, False: no, use images)
    if not strToBool(sys.argv[2]):
        produce_output()
    tm.makeTree()

def eraseFile():
    open(valDir, "w").close()

def getHierarchy(plantName, driver):
    sbox = driver.find_element(By.XPATH, "//input[@type='text' and @autofocus='autofocus' and @name='string']")
    sbox.send_keys(plantName)

    submit = driver.find_element(By.XPATH, "//input[@type='submit' and @class='clbutton' and @value=' OK ']")
    submit.click()

    try:
        driver.find_element(By.XPATH, '//*[@id="screen"]/div[3]/div/p/span[1]')
    except:
        try:
            firstElem = driver.find_element(By.XPATH, '//*[@id="screen"]/div[5]/div[1]/div/a')
            firstElem.click()
        except NoSuchElementException:
            try:
                secElem = driver.find_element(By.XPATH, '//*[@id="screen"]/div[6]/div[1]/div/a')
                secElem.click()
            except:
                print(f"Error! Could not find {plantName}")
                return

    # page_source = driver.page_source
    # soup = BS(page_source, features="lxml")

    #Scrape all values
    try:
        try:
            kingdom = getTypeElement(driver, '//*[@id="screen"]/div[3]/div/p/span[1]/a', '//*[@id="screen"]/div[3]/div/p/span[1]/b')
        except:
            kingdom = "-"
        try:
            phylum = getTypeElement(driver, '//*[@id="screen"]/div[3]/div/p/span[2]/a', '//*[@id="screen"]/div[3]/div/p/span[2]/b')
        except:
            phylum = "-"
        try:
            plantClass = getTypeElement(driver, '//*[@id="screen"]/div[3]/div/p/span[3]/a', '//*[@id="screen"]/div[3]/div/p/span[3]/b')
        except:
            plantClass = "-"
        try:
            order = getTypeElement(driver, '//*[@id="screen"]/div[3]/div/p/span[4]/a', '//*[@id="screen"]/div[3]/div/p/span[4]/b')
        except:
            order = "-"
        try:
            family = getTypeElement(driver, '//*[@id="screen"]/div[3]/div/p/span[5]/a', '//*[@id="screen"]/div[3]/div/p/span[5]/b')
        except:
            family = "-"
        try:
            genus = getTypeElement(driver, '//*[@id="screen"]/div[3]/div/p/span[6]/a', '//*[@id="screen"]/div[3]/div/p/span[6]/b')
        except:
            genus = "-"
    except Exception as e:
        print(f"Failed to find: {plantName}")
        return

    try:
        czechName = driver.find_element(By.XPATH, '//*[@id="screen"]/div[3]/div/h1/strong[1]').text
    except:
        czechName = "-"
    try:
        latinName = driver.find_element(By.XPATH, '//*[@id="screen"]/div[3]/div/h1/strong[2]/em').text
    except:
        try:
            latinName = driver.find_element(By.XPATH, '//*[@id="screen"]/div[3]/div/h1/strong/em').text
        except:
            latinName = "-"

    result = plant(kingdom, phylum, plantClass, order, family, genus, czechName, latinName)

    #Create organized output
    familyDict = {family: [f"{genus}NAME:{czechName} \n({latinName})"]}
    orderDict = {order: familyDict}
    classDict = {plantClass: orderDict}
    phylumDict = {phylum: classDict}
    kingdomDict = {kingdom: phylumDict}
    
    if kingdom in iterableKeys(resultDict):
        if phylum in iterableKeys(resultDict[kingdom]):
            if plantClass in iterableKeys(resultDict[kingdom][phylum]):
                if order in iterableKeys(resultDict[kingdom][phylum][plantClass]):
                    if family in iterableKeys(resultDict[kingdom][phylum][plantClass][order]):
                        if genus in resultDict[kingdom][phylum][plantClass][order][family]:
                            return result
                        else:
                            if type(resultDict[kingdom][phylum][plantClass][order][family]) == list:
                                resultDict[kingdom][phylum][plantClass][order][family].append(f"{genus}NAME:{czechName} \n({latinName})")
                            else:
                                resultDict[kingdom][phylum][plantClass][order][family] = [f"{genus}NAME:{czechName} \n({latinName})"]
                    else:
                        resultDict[kingdom][phylum][plantClass][order][family] = [f"{genus}NAME:{czechName} \n({latinName})"]
                else:
                    resultDict[kingdom][phylum][plantClass][order] = familyDict
            else:
                resultDict[kingdom][phylum][plantClass] = orderDict
        else:
            resultDict[kingdom][phylum] = classDict
    else:
        resultDict[kingdom] = phylumDict

    return result

def iterableKeys(diction):
    resList = []
    for key, value in diction.items():
        resList.append(key)
    return resList

def strToBool(inp):
    return inp in ["True", "true"]

def merge_dicts(dict_list):
    result = {}
    for sub_dict in dict_list:
        for key, value in sub_dict.items():
            if isinstance(value, dict):
                result[key] = merge_dicts([result.get(key, {}), value])
            elif key not in result:
                result[key] = value
            elif value != result[key]:
                result[key] = [result[key], value]
    return result

def getElementContent(soup, XPath):
    identifier = soup.find()

def getTypeElement(driver, latinPath, czechPath):
    # identifier = soup.find(string=re.compile(input))
    # if input == "rod" and identifier.find_next("a") == "plantae":
    #     print(type(identifier))
    # latin = identifier.find_next("a")
    # czech = latin.find_next("b")
    try:
        czech = driver.find_element(By.XPATH, czechPath).text
    except Exception as e:
        # print(f"{czechPath}: {e}")
        czech = "-"
    try:
        latin = driver.find_element(By.XPATH, latinPath).text
    except Exception as e:
        # print(f"{latinPath}: {e}")
        latin = "-"
    
    # if len(latin) > 0:
    #     latinText = str(latin.contents[0])
    #     if latinText[:4] == "<em>":
    #         latinText = latinText[4:-5]
    # else:
    #     latinText = "-"

    # czechText = "-"
    # if len(czech) > 0:
    #     czechText = czech.contents[0]

    # if input == "rod" and czechText == "rostliny":
    #     return getTypeElement(soup, "podčeleď")
    return f"{czech} ({latin})"

if __name__ == "__main__":
    # app.run(host="127.0.0.1", port=8080, debug=True)
    main()